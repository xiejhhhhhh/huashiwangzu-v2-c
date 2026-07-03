from app.core.exceptions import ValidationError
from app.middleware.auth import require_permission
from app.models.user import User
from app.schemas.common import ApiResponse
from app.services.module_registry import register_capability
from app.services.parser_resource_diagnostics import (
    build_resource_diagnostic,
    store_extracted_resources_with_diagnostics,
)
from app.services.uploaded_file_runner import run_uploaded_file_capability
from fastapi import APIRouter, Depends
from pydantic import BaseModel, Field

router = APIRouter(prefix="/api/pptx-parser", tags=["pptx-parser"])


class ParseRequest(BaseModel):
    file_id: int = Field(gt=0)


async def _parse(params: dict, caller: str) -> dict:
    import base64
    import hashlib

    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    allowed = {"pptx"}
    file_id = params.get("file_id")
    if not isinstance(file_id, int) or file_id <= 0:
        raise ValidationError("file_id must be a positive integer")

    def parse_file(file_id, _file, full_path, _ext):
        try:
            prs = Presentation(str(full_path))
        except Exception as exc:
            raise ValidationError(f"Failed to parse PPTX file: {exc}") from exc

        blocks = []
        resources = []
        resource_diagnostics = []
        resource_counter = 0

        for slide_idx, slide in enumerate(prs.slides):
            pno = slide_idx + 1
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue
                        block_type = "heading" if ("title" in str(shape.name).lower() or "标题" in str(shape.name)) else "paragraph"
                        blocks.append({"type": block_type, "text": text, "page": pno, "resource_ref": None})
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    resource_counter += 1
                    mime_type = "image/png"
                    extract_diagnostic_recorded = False
                    try:
                        img = shape.image
                        img_bytes = img.blob
                        mime_type = getattr(img, "content_type", None) or "image/png"
                    except Exception as exc:
                        extract_diagnostic_recorded = True
                        img_bytes = b""
                        resource_diagnostics.append(build_resource_diagnostic(
                            parser="pptx-parser",
                            stage="extract",
                            status="degraded",
                            code="resource_extract_failed",
                            message="Failed to extract PPTX embedded image bytes.",
                            resource={
                                "id": resource_counter,
                                "type": "image",
                                "page": pno,
                                "mime_type": mime_type,
                                "filename": f"slide{pno}_{hashlib.md5(str(shape.name).encode()).hexdigest()[:8]}.png",
                                "description": f"Slide {pno} image ({shape.name})",
                            },
                            error=exc,
                        ))

                    blocks.append({"type": "image", "text": "", "page": pno, "resource_ref": resource_counter})
                    resources.append({
                        "id": resource_counter,
                        "type": "image",
                        "page": pno,
                        "mime_type": mime_type,
                        "filename": f"slide{pno}_{hashlib.md5(str(shape.name).encode()).hexdigest()[:8]}.png",
                        "description": f"Slide {pno} image ({shape.name})",
                        "_resource_diagnostic_recorded": extract_diagnostic_recorded,
                        "_bytes_b64": base64.b64encode(img_bytes).decode("ascii") if img_bytes else "",
                    })

        return {
            "file_id": file_id,
            "format": "pptx",
            "blocks": blocks,
            "resources": resources,
            "resource_diagnostics": resource_diagnostics,
        }

    result = await run_uploaded_file_capability(params, caller, allowed, parse_file)
    return await store_extracted_resources_with_diagnostics(result, caller=caller, parser="pptx-parser")


@router.get("/health")
async def health():
    return ApiResponse(data={"module": "pptx-parser", "status": "ok"})


@router.post("/parse")
async def call_parse(payload: ParseRequest, user: User = Depends(require_permission("viewer"))):
    result = await _parse({"file_id": payload.file_id}, f"user:{user.id}")
    return ApiResponse(data=result)


register_capability(
    "pptx-parser", "parse", _parse,
    description="Parse PPTX files into unified content blocks",
    brief="解析 PPTX 文档",
    parameters={"file_id": {"type": "int"}},
    min_role="viewer",
)
