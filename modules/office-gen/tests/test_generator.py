"""Tests for office-gen: verify generated files are valid and non-empty.

Run from project root:
  cd backend && .venv/bin/python -m pytest ../modules/office-gen/tests/test_generator.py -v
"""
from __future__ import annotations

import importlib.util
import io
import sys
import types
from pathlib import Path

import pytest

# ── Build package hierarchy for "modules/office-gen/backend/" ──
# Since the directory has a hyphen, we use the normalized name "office_gen"

_PROJECT_ROOT = Path(__file__).resolve().parents[3]


def _load_backend_module(name: str) -> types.ModuleType:
    """Load a modules/office-gen/backend/*.py module with proper package prefix."""
    full_name = f"office_gen.{name}"
    file_path = _PROJECT_ROOT / "modules" / "office-gen" / "backend" / f"{name}.py"
    if not file_path.exists():
        raise RuntimeError(f"Module file not found: {file_path}")
    spec = importlib.util.spec_from_file_location(full_name, file_path)
    if not spec or not spec.loader:
        raise RuntimeError(f"Cannot load module: {file_path}")
    module = importlib.util.module_from_spec(spec)
    sys.modules[full_name] = module
    spec.loader.exec_module(module)
    return module


# Load the generator module
gen = _load_backend_module("generator")


def _check_docx(data: bytes) -> bool:
    try:
        from docx import Document
        Document(io.BytesIO(data))
        return True
    except Exception:
        return False


def _check_xlsx(data: bytes) -> bool:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(data))
        return len(wb.sheetnames) > 0
    except Exception:
        return False


def _check_pptx(data: bytes) -> bool:
    try:
        from pptx import Presentation
        Presentation(io.BytesIO(data))
        return True
    except Exception:
        return False


def _check_pdf(data: bytes) -> bool:
    return data.startswith(b"%PDF") and data.rstrip().endswith(b"%%EOF")


class TestDocxGenerator:
    def test_basic_document(self):
        data = gen.generate_docx({
            "filename": "test",
            "content": [
                {"type": "标题", "text": "Hello World", "level": 1},
                {"type": "段落", "text": "This is a test paragraph", "bold": True},
            ],
        })
        assert len(data) > 100
        assert _check_docx(data)

    def test_with_table(self):
        data = gen.generate_docx({
            "filename": "test_table",
            "content": [
                {"type": "标题", "text": "Table Test", "level": 2},
                {"type": "表格", "header": ["Name", "Age", "City"],
                 "rows": [["Zhang", "28", "Beijing"], ["Li", "32", "Shanghai"]]},
            ],
        })
        assert _check_docx(data)

    def test_empty_content(self):
        with pytest.raises(ValueError, match="content must be a non-empty array"):
            gen.generate_docx({"filename": "empty", "content": []})

    def test_alignment(self):
        data = gen.generate_docx({
            "filename": "align",
            "content": [
                {"type": "段落", "text": "Left", "align": "left"},
                {"type": "段落", "text": "Center", "align": "center"},
            ],
        })
        assert _check_docx(data)

    def test_content_ir_english_blocks(self):
        data = gen.generate_docx({
            "filename": "content_ir",
            "content": [
                {"type": "heading", "text": "Content IR Title", "data": {"level": 1}},
                {"type": "paragraph", "text": "Paragraph from Content IR"},
                {"type": "table", "table_header": ["Name", "Count"], "table_rows": [["Alpha", 0]]},
            ],
        })
        doc = __import__("docx").Document(io.BytesIO(data))
        assert "Content IR Title" in [p.text for p in doc.paragraphs]
        assert "Paragraph from Content IR" in [p.text for p in doc.paragraphs]
        assert doc.tables[0].cell(1, 1).text == "0"

    def test_content_ir_blocks_alias_and_data_fields(self):
        data = gen.generate_docx({
            "filename": "content_ir_alias",
            "blocks": [
                {"type": "heading", "data": {"text": "Alias Title", "level": 2}},
                {"type": "table", "data": {"headers": ["Name", "Count"], "rows": [{"Name": "Beta", "Count": 0}]}},
            ],
        })
        doc = __import__("docx").Document(io.BytesIO(data))
        assert "Alias Title" in [p.text for p in doc.paragraphs]
        assert doc.tables[0].cell(1, 1).text == "0"


class TestXlsxGenerator:
    def test_basic_spreadsheet(self):
        data = gen.generate_xlsx({
            "filename": "test",
            "sheets": [{
                "name": "Sheet1",
                "columns": ["A", "B", "C"],
                "rows": [[1, 2, 3], [4, 5, 6], [7, 8, 9]],
            }],
        })
        assert len(data) > 100
        assert _check_xlsx(data)

    def test_multiple_sheets(self):
        data = gen.generate_xlsx({
            "filename": "multi",
            "sheets": [
                {"name": "Page1", "columns": ["X"], "rows": [[1]]},
                {"name": "Page2", "columns": ["Y"], "rows": [[2]]},
            ],
        })
        assert _check_xlsx(data)
        wb = __import__("openpyxl").load_workbook(io.BytesIO(data))
        assert wb.sheetnames == ["Page1", "Page2"]

    def test_empty_sheets(self):
        with pytest.raises(ValueError, match="sheets must be a non-empty array"):
            gen.generate_xlsx({"filename": "empty", "sheets": []})

    def test_content_ir_column_dicts_and_row_dicts(self):
        data = gen.generate_xlsx({
            "filename": "ir_sheet",
            "sheets": [{
                "name": "Sheet1",
                "columns": [{"name": "item"}, {"name": "count"}],
                "rows": [{"item": "Alpha", "count": 0}],
            }],
        })
        wb = __import__("openpyxl").load_workbook(io.BytesIO(data))
        ws = wb["Sheet1"]
        assert ws["A1"].value == "item"
        assert ws["B2"].value == "0"

    def test_content_ir_sheet_blocks_with_child_table(self):
        data = gen.generate_xlsx({
            "filename": "ir_sheet_blocks",
            "content_ir": {
                "content_type": "spreadsheet",
                "blocks": [{
                    "type": "sheet",
                    "data": {"name": "Data"},
                    "children": [{
                        "type": "table",
                        "data": {
                            "headers": [{"name": "item"}, {"name": "count"}],
                            "rows": [{"item": "Beta", "count": 0}],
                        },
                    }],
                }],
            },
        })
        wb = __import__("openpyxl").load_workbook(io.BytesIO(data))
        ws = wb["Data"]
        assert ws["A1"].value == "item"
        assert ws["B2"].value == "0"


class TestPptxGenerator:
    def test_basic_presentation(self):
        data = gen.generate_pptx({
            "filename": "test",
            "slides": [
                {"title": "Slide 1", "bullets": ["Point 1", "Point 2"], "notes": "Notes text"},
                {"title": "Slide 2", "bullets": [{"text": "Sub item", "level": 1}]},
            ],
        })
        assert len(data) > 100
        assert _check_pptx(data)

    def test_empty_slides(self):
        with pytest.raises(ValueError, match="slides must be a non-empty array"):
            gen.generate_pptx({"filename": "empty", "slides": []})

    def test_content_ir_slide_elements(self):
        data = gen.generate_pptx({
            "filename": "ir_deck",
            "slides": [{
                "name": "Content IR Slide",
                "elements": [
                    {"type": "heading", "text": "First point"},
                    {"type": "paragraph", "text": "Second point", "level": 1},
                ],
            }],
        })
        prs = __import__("pptx").Presentation(io.BytesIO(data))
        assert prs.slides[0].shapes.title.text == "Content IR Slide"
        slide_text = "\n".join(shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text"))
        assert "First point" in slide_text

    def test_content_ir_slide_blocks_with_children(self):
        data = gen.generate_pptx({
            "filename": "ir_deck_blocks",
            "content_ir": {
                "content_type": "presentation",
                "blocks": [{
                    "type": "slide",
                    "data": {"title": "IR Child Slide"},
                    "children": [
                        {"type": "heading", "data": {"text": "Child heading"}},
                        {"type": "paragraph", "text": "Child paragraph", "level": 1},
                    ],
                }],
            },
        })
        prs = __import__("pptx").Presentation(io.BytesIO(data))
        assert prs.slides[0].shapes.title.text == "IR Child Slide"
        slide_text = "\n".join(shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text"))
        assert "Child heading" in slide_text


class TestPdfGenerator:
    def test_basic_pdf(self):
        data = gen.generate_pdf({
            "filename": "test",
            "content": [
                {"type": "标题", "text": "Test Document", "level": 1},
                {"type": "段落", "text": "This is a test paragraph."},
            ],
        })
        assert len(data) > 500
        assert _check_pdf(data)

    def test_with_table_pdf(self):
        data = gen.generate_pdf({
            "filename": "test_table",
            "content": [
                {"type": "标题", "text": "Table Test", "level": 2},
                {"type": "表格", "header": ["Name", "Value"],
                 "rows": [["Item A", "100"], ["Item B", "200"]]},
            ],
        })
        assert _check_pdf(data)

    def test_empty_content_pdf(self):
        with pytest.raises(ValueError, match="content must be a non-empty array"):
            gen.generate_pdf({"filename": "empty", "content": []})

    def test_content_ir_english_pdf_blocks(self):
        data = gen.generate_pdf({
            "filename": "ir_pdf",
            "content": [
                {"type": "heading", "text": "PDF IR Title", "data": {"level": 1}},
                {"type": "paragraph", "text": "PDF paragraph"},
                {"type": "table", "header": ["Name", "Count"], "rows": [["Alpha", 0]]},
            ],
        })
        assert len(data) > 500
        assert _check_pdf(data)

    def test_content_ir_pdf_blocks_alias_and_data_table(self):
        data = gen.generate_pdf({
            "filename": "ir_pdf_alias",
            "blocks": [
                {"type": "heading", "data": {"text": "PDF Alias Title", "level": 1}},
                {"type": "table", "data": {"headers": ["Name", "Count"], "rows": [{"Name": "Beta", "Count": 0}]}},
            ],
        })
        assert len(data) > 500
        assert _check_pdf(data)


class TestConverter:
    @staticmethod
    def _get_converter():
        conv = _load_backend_module("converter")
        return conv

    def test_check_libreoffice(self):
        conv = self._get_converter()
        result = conv.check_libreoffice()
        assert result is None or isinstance(result, str)

    def test_get_install_instructions(self):
        conv = self._get_converter()
        instructions = conv.get_install_instructions()
        assert "LibreOffice" in instructions
        assert "brew install" in instructions
