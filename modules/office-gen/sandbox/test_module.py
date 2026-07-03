"""Sandbox validation for office-gen.

Runs without DB writes: loads the module generator directly and verifies that
docx/xlsx/pptx/pdf outputs are real files, including Content IR-style English
blocks used by the framework export pipeline.
"""
from __future__ import annotations

import asyncio
import importlib.util
import io
import os
import sys
import types
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parents[3]
BACKEND_ROOT = PROJECT_ROOT / "backend"
os.environ.setdefault("JWT_SECRET", "office-gen-sandbox-test-secret")
if str(BACKEND_ROOT) not in sys.path:
    sys.path.insert(0, str(BACKEND_ROOT))


def _load_backend_module(name: str) -> types.ModuleType:
    package_name = "office_gen_sandbox"
    backend_dir = PROJECT_ROOT / "modules" / "office-gen" / "backend"
    if package_name not in sys.modules:
        package = types.ModuleType(package_name)
        setattr(package, "__path__", [str(backend_dir)])
        sys.modules[package_name] = package

    full_name = f"{package_name}.{name}"
    file_path = backend_dir / f"{name}.py"
    spec = importlib.util.spec_from_file_location(full_name, file_path)
    if not spec or not spec.loader:
        raise RuntimeError(f"Cannot load module: {file_path}")
    module = importlib.util.module_from_spec(spec)
    sys.modules[full_name] = module
    spec.loader.exec_module(module)
    return module


gen = _load_backend_module("generator")
router = _load_backend_module("router")


def _assert_docx(data: bytes) -> None:
    from docx import Document

    assert len(data) > 100
    doc = Document(io.BytesIO(data))
    text = "\n".join(p.text for p in doc.paragraphs)
    assert "Content IR Title" in text
    assert "Content IR paragraph" in text


def _assert_xlsx(data: bytes) -> None:
    from openpyxl import load_workbook

    assert len(data) > 100
    wb = load_workbook(io.BytesIO(data))
    ws = wb["Sheet1"]
    assert ws["A1"].value == "item"
    assert ws["B2"].value == "0"


def _assert_pptx(data: bytes) -> None:
    from pptx import Presentation

    assert len(data) > 100
    prs = Presentation(io.BytesIO(data))
    assert prs.slides[0].shapes.title.text == "Content IR Slide"
    slide_text = "\n".join(shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text"))
    assert "First point" in slide_text


def _assert_pdf(data: bytes) -> None:
    assert len(data) > 500
    assert data.startswith(b"%PDF")
    assert data.rstrip().endswith(b"%%EOF")


def test_docx_generation_accepts_content_ir_blocks() -> None:
    data = gen.generate_docx({
        "filename": "sandbox_doc",
        "content": [
            {"type": "heading", "text": "Content IR Title", "data": {"level": 1}},
            {"type": "paragraph", "text": "Content IR paragraph"},
            {"type": "table", "header": ["item", "count"], "rows": [["Alpha", 0]]},
        ],
    })
    _assert_docx(data)


def test_xlsx_generation_accepts_content_ir_tables() -> None:
    data = gen.generate_xlsx({
        "filename": "sandbox_sheet",
        "content_ir": {
            "content_type": "spreadsheet",
            "blocks": [{
                "type": "sheet",
                "data": {"name": "Sheet1"},
                "children": [{
                    "type": "table",
                    "data": {
                        "headers": [{"name": "item"}, {"name": "count"}],
                        "rows": [{"item": "Alpha", "count": 0}],
                    },
                }],
            }],
        },
    })
    _assert_xlsx(data)


def test_pptx_generation_accepts_content_ir_elements() -> None:
    data = gen.generate_pptx({
        "filename": "sandbox_deck",
        "content_ir": {
            "content_type": "presentation",
            "blocks": [{
                "type": "slide",
                "data": {"title": "Content IR Slide"},
                "children": [
                    {"type": "heading", "text": "First point"},
                    {"type": "paragraph", "text": "Second point", "level": 1},
                ],
            }],
        },
    })
    _assert_pptx(data)


def test_pdf_generation_accepts_content_ir_blocks() -> None:
    data = gen.generate_pdf({
        "filename": "sandbox_pdf",
        "blocks": [
            {"type": "heading", "data": {"text": "PDF Title", "level": 1}},
            {"type": "paragraph", "data": {"text": "PDF paragraph"}},
            {"type": "table", "data": {"headers": ["item", "count"], "rows": [{"item": "Alpha", "count": 0}]}},
        ],
    })
    _assert_pdf(data)


def test_generators_reject_empty_outputs() -> None:
    checks = [
        (gen.generate_docx, {"filename": "empty", "content": []}),
        (gen.generate_xlsx, {"filename": "empty", "sheets": []}),
        (gen.generate_pptx, {"filename": "empty", "slides": []}),
        (gen.generate_pdf, {"filename": "empty", "content": []}),
    ]
    for func, payload in checks:
        try:
            func(payload)
        except ValueError:
            continue
        raise AssertionError(f"{func.__name__} accepted empty payload")


def test_router_http_models_accept_content_ir_aliases() -> None:
    docx = router.GenerateDocxRequest(
        filename="alias_doc",
        blocks=[{"type": "paragraph", "text": "ok"}],
    ).model_dump()
    assert docx["blocks"][0]["text"] == "ok"

    xlsx = router.GenerateXlsxRequest(
        filename="alias_sheet",
        content_ir={"blocks": [{"type": "sheet", "data": {"name": "Data"}}]},
    ).model_dump()
    assert xlsx["content_ir"]["blocks"][0]["type"] == "sheet"


def test_router_helpers_reject_bad_parameters() -> None:
    from app.core.exceptions import ValidationError

    for value in (None, "", 0, -1, "abc", True):
        try:
            router._positive_int(value, "file_id")
        except ValidationError:
            continue
        raise AssertionError(f"_positive_int accepted bad value: {value!r}")

    for target_format in (None, "", "exe"):
        try:
            router._normalize_convert_format(target_format)
        except ValidationError:
            continue
        raise AssertionError(f"_normalize_convert_format accepted bad value: {target_format!r}")

    try:
        router._render_generated_bytes(lambda: (_ for _ in ()).throw(RuntimeError("missing dependency")))
    except ValidationError:
        pass
    else:
        raise AssertionError("_render_generated_bytes accepted RuntimeError")


def test_capability_bad_parameters_fail_before_io() -> None:
    from app.core.exceptions import ValidationError

    async def run_checks() -> None:
        cases = [
            router._cap_convert({"file_id": "abc"}, "user:1"),
            router._cap_convert({"file_id": 0}, "user:1"),
            router._cap_replace_existing({"format": "xlsx", "target_file_id": "abc"}, "user:1"),
            router._cap_export_to_artifact({"file_id": False}, "user:1"),
        ]
        for case in cases:
            try:
                await case
            except ValidationError:
                continue
            raise AssertionError("bad capability parameters did not raise ValidationError")

    asyncio.run(run_checks())


def main() -> None:
    tests = [
        test_docx_generation_accepts_content_ir_blocks,
        test_xlsx_generation_accepts_content_ir_tables,
        test_pptx_generation_accepts_content_ir_elements,
        test_pdf_generation_accepts_content_ir_blocks,
        test_generators_reject_empty_outputs,
        test_router_http_models_accept_content_ir_aliases,
        test_router_helpers_reject_bad_parameters,
        test_capability_bad_parameters_fail_before_io,
    ]
    print("=" * 60)
    print("office-gen sandbox test")
    print("=" * 60)
    for test in tests:
        test()
        print(f"  [PASS] {test.__name__}")
    print("=" * 60)
    print("PASS: office-gen sandbox test")


if __name__ == "__main__":
    main()
