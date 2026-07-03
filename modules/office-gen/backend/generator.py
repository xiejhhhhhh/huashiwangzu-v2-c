"""Office document generator — JSON intermediate layer → file bytes.

Each function takes a validated JSON dict and returns (bytes, mime_type).
"""
import io
import logging

logger = logging.getLogger("v2.office-gen").getChild("generator")

try:
    from docx.shared import Cm
except ImportError:
    Cm = None


# ── DOCX generator ──────────────────────────────────────────────────────

def _block_type(block: dict, default: str = "段落") -> str:
    return str(block.get("类型", block.get("type", default)) or default).lower()


def _block_text(block: dict) -> str:
    data = block.get("data") if isinstance(block.get("data"), dict) else {}
    for key in ("文本", "text", "title", "name", "value"):
        if key in block and block[key] is not None:
            return str(block[key])
    for key in ("text", "title", "name", "value"):
        if key in data and data[key] is not None:
            return str(data[key])
    return ""


def _block_level(block: dict) -> int:
    raw = block.get("级别", block.get("level", block.get("data", {}).get("level", 1)))
    try:
        return int(raw or 1)
    except (TypeError, ValueError):
        return 1


def _block_align(block: dict) -> str:
    return str(block.get("对齐", block.get("align", "left")) or "left").lower()


def _block_bold(block: dict) -> bool:
    return bool(block.get("加粗", block.get("bold", False)))


def _is_heading(block_type: str) -> bool:
    return block_type in {"标题", "heading", "head", "title", "h1", "h2", "h3", "h4"}


def _is_paragraph(block_type: str) -> bool:
    return block_type in {"段落", "paragraph", "text", "textbox", "list", "code"}


def _is_table(block_type: str) -> bool:
    return block_type in {"表格", "table"}


def _is_image(block_type: str) -> bool:
    return block_type in {"图片", "image"}


def _is_page_break(block_type: str) -> bool:
    return block_type in {"分页", "page_break", "pagebreak"}


def _table_header(block: dict) -> list:
    data = block.get("data") if isinstance(block.get("data"), dict) else {}
    return (
        block.get("表头")
        or block.get("header")
        or block.get("table_header")
        or data.get("headers")
        or data.get("columns")
        or data.get("header")
        or []
    )


def _table_rows(block: dict) -> list:
    data = block.get("data") if isinstance(block.get("data"), dict) else {}
    return block.get("行") or block.get("rows") or block.get("table_rows") or data.get("rows") or []


def _as_list(value) -> list:
    return value if isinstance(value, list) else []


def _content_blocks(params: dict) -> list:
    blocks = params.get("content") or params.get("blocks")
    if not blocks and isinstance(params.get("content_ir"), dict):
        blocks = params["content_ir"].get("blocks")
    return _as_list(blocks)


def _normalize_row(row, columns: list | None = None) -> list:
    if isinstance(row, dict):
        if columns:
            return [_cell_value(row.get(_cell_value(column))) for column in columns]
        return [_cell_value(value) for value in row.values()]
    if isinstance(row, (list, tuple)):
        return [_cell_value(value) for value in row]
    return [_cell_value(row)]


def _normalize_table_data(block: dict) -> list[list[str]]:
    header = _table_header(block)
    rows = _table_rows(block)
    table_rows: list[list[str]] = []
    normalized_header = [_cell_value(cell) for cell in header] if header else []
    if normalized_header:
        table_rows.append(normalized_header)
    for row in rows:
        table_rows.append(_normalize_row(row, header))
    return [row for row in table_rows if row]


def _require_non_empty(name: str, items: list) -> None:
    if not items:
        raise ValueError(f"{name} must be a non-empty array")


def _cell_value(value) -> str:
    if isinstance(value, dict):
        for key in ("name", "label", "text", "value", "key"):
            if key in value:
                return "" if value[key] is None else str(value[key])
        return str(value)
    return "" if value is None else str(value)


def generate_docx(params: dict) -> bytes:
    try:
        from docx import Document
        from docx.shared import Pt
    except ImportError:
        raise RuntimeError("python-docx is not installed. Run: pip install python-docx")

    doc = Document()
    content = _content_blocks(params)
    _require_non_empty("content", content)
    rendered_blocks = 0

    for block in content:
        block_type = _block_type(block)
        text = _block_text(block)
        bold = _block_bold(block)
        align = _block_align(block)
        level = _block_level(block)

        if _is_heading(block_type):
            doc.add_heading(text, level=min(level, 4))
            rendered_blocks += 1
        elif _is_paragraph(block_type):
            p = doc.add_paragraph()
            run = p.add_run(text)
            run.bold = bold
            run.font.size = Pt(12)
            _set_alignment(p, align)
            rendered_blocks += 1
        elif _is_table(block_type):
            rows = _normalize_table_data(block)
            if rows:
                col_count = max(len(row) for row in rows)
                table = doc.add_table(rows=len(rows), cols=col_count)
                table.style = "Light Grid Accent 1"
                for i, row_data in enumerate(rows):
                    for j, cell_text in enumerate(row_data):
                        if j < len(table.rows[i].cells):
                            table.rows[i].cells[j].text = _cell_value(cell_text)
                rendered_blocks += 1
        elif _is_image(block_type):
            if _add_image_block(doc, block):
                rendered_blocks += 1
        elif _is_page_break(block_type):
            doc.add_page_break()

    if rendered_blocks == 0:
        raise ValueError("content must contain at least one renderable block")

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()


def _set_alignment(paragraph, align: str):
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    mapping = {"left": WD_ALIGN_PARAGRAPH.LEFT, "center": WD_ALIGN_PARAGRAPH.CENTER,
               "right": WD_ALIGN_PARAGRAPH.RIGHT, "justify": WD_ALIGN_PARAGRAPH.JUSTIFY}
    paragraph.alignment = mapping.get(align, WD_ALIGN_PARAGRAPH.LEFT)


def _add_image_block(doc, block: dict) -> bool:
    if Cm is None:
        raise RuntimeError("python-docx is not installed. Run: pip install python-docx")

    image_data = block.get("image_data")
    image_path = block.get("image_path")
    width = block.get("width", 14)
    height = block.get("height")

    try:
        if image_data and isinstance(image_data, str):
            import base64
            raw = base64.b64decode(image_data)
            buf = io.BytesIO(raw)
            if height:
                doc.add_picture(buf, width=Cm(width), height=Cm(height))
            else:
                doc.add_picture(buf, width=Cm(width))
            return True
        elif image_path:
            doc.add_picture(image_path, width=Cm(width))
            return True
    except Exception as exc:
        logger.warning("Failed to add image: %s", exc)
    return False


# ── XLSX generator ──────────────────────────────────────────────────────

def generate_xlsx(params: dict) -> bytes:
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Alignment, Font, PatternFill
    except ImportError:
        raise RuntimeError("openpyxl is not installed. Run: pip install openpyxl")

    wb = Workbook()
    sheets = _xlsx_sheets(params)
    _require_non_empty("sheets", sheets)
    wb.remove(wb.active)
    written_cells = 0
    for sheet_spec in sheets:
        ws = wb.create_sheet(title=_sheet_name(sheet_spec))
        columns, rows = _sheet_table_data(sheet_spec)

        if columns:
            ws.append([_cell_value(column) for column in columns])
            written_cells += len(columns)
            for cell in ws[1]:
                cell.font = Font(bold=True, color="FFFFFF")
                cell.fill = PatternFill(start_color="2395BC", end_color="2395BC", fill_type="solid")
                cell.alignment = Alignment(horizontal="center")

        for row_data in rows:
            if isinstance(row_data, dict):
                row_values = [_cell_value(row_data.get(_cell_value(column))) for column in columns] if columns else [
                    _cell_value(value) for value in row_data.values()
                ]
            else:
                row_values = [_cell_value(value) for value in list(row_data)]
            ws.append(row_values)
            written_cells += len(row_values)

        for col in ws.columns:
            max_len = max((len(str(c.value or "")) for c in col), default=8)
            ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 60)

    if written_cells == 0:
        raise ValueError("sheets must contain at least one row, column, or cell")

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()


def _xlsx_sheets(params: dict) -> list:
    sheets = params.get("工作表") or params.get("sheets")
    if not sheets and isinstance(params.get("content_ir"), dict):
        sheets = params["content_ir"].get("blocks")
    if not sheets:
        sheets = params.get("blocks")
    return _as_list(sheets)


def _sheet_name(sheet_spec: dict) -> str:
    data = sheet_spec.get("data") if isinstance(sheet_spec.get("data"), dict) else {}
    name = (
        sheet_spec.get("表名")
        or sheet_spec.get("name")
        or sheet_spec.get("title")
        or sheet_spec.get("text")
        or data.get("name")
        or data.get("title")
        or "Sheet"
    )
    return str(name)[:31] or "Sheet"


def _sheet_table_data(sheet_spec: dict) -> tuple[list, list]:
    data = sheet_spec.get("data") if isinstance(sheet_spec.get("data"), dict) else {}
    columns = sheet_spec.get("列") or sheet_spec.get("columns") or data.get("columns") or data.get("headers") or []
    rows = sheet_spec.get("行") or sheet_spec.get("rows") or data.get("rows") or []
    if columns or rows:
        return _as_list(columns), _as_list(rows)

    children = _as_list(sheet_spec.get("children"))
    table_blocks = [child for child in children if isinstance(child, dict) and _is_table(_block_type(child))]
    if table_blocks:
        first = table_blocks[0]
        return _table_header(first), _table_rows(first)
    range_blocks = [child for child in children if isinstance(child, dict) and _block_type(child) == "range"]
    if range_blocks:
        range_data = range_blocks[0].get("data") if isinstance(range_blocks[0].get("data"), dict) else {}
        return _as_list(range_data.get("columns") or range_data.get("headers")), _as_list(range_data.get("rows") or range_data.get("values"))
    return [], []


# ── PPTX generator ──────────────────────────────────────────────────────

def generate_pptx(params: dict) -> bytes:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        raise RuntimeError("python-pptx is not installed. Run: pip install python-pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = _pptx_slides(params)
    _require_non_empty("slides", slides)
    rendered_slides = 0
    for slide_spec in slides:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        title = slide.shapes.title
        title_text = _slide_title(slide_spec)
        if title:
            title.text = title_text

        bullets = _slide_bullets(slide_spec)
        if bullets:
            body = slide.shapes.placeholders[1]
            tf = body.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                if isinstance(bullet, str):
                    p.text = bullet
                elif isinstance(bullet, dict):
                    p.text = bullet.get("text", "")
                    p.level = bullet.get("level", 0)

        notes_text = slide_spec.get("备注", slide_spec.get("notes", ""))
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text
        if title_text or bullets or notes_text:
            rendered_slides += 1

    if rendered_slides == 0:
        raise ValueError("slides must contain at least one renderable slide")

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


def _pptx_slides(params: dict) -> list:
    slides = params.get("幻灯片") or params.get("slides")
    if not slides and isinstance(params.get("content_ir"), dict):
        slides = params["content_ir"].get("blocks")
    if not slides:
        slides = params.get("blocks")
    return _as_list(slides)


def _slide_title(slide_spec: dict) -> str:
    data = slide_spec.get("data") if isinstance(slide_spec.get("data"), dict) else {}
    return str(
        slide_spec.get("标题")
        or slide_spec.get("title")
        or slide_spec.get("name")
        or slide_spec.get("text")
        or data.get("title")
        or data.get("name")
        or ""
    )


def _slide_bullets(slide_spec: dict) -> list:
    bullets = slide_spec.get("要点") or slide_spec.get("bullets")
    if bullets:
        return _as_list(bullets)
    elements = slide_spec.get("elements")
    if not elements:
        elements = slide_spec.get("children")
    if not isinstance(elements, list):
        return []
    return [
        {"text": _block_text(elem), "level": _block_level(elem)}
        for elem in elements
        if isinstance(elem, dict) and _block_text(elem)
    ]


# ── PDF generator ───────────────────────────────────────────────────────

def generate_pdf(params: dict) -> bytes:
    try:
        from reportlab.lib.colors import HexColor
        from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT, TA_RIGHT
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import cm
        from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
    except ImportError:
        raise RuntimeError("reportlab is not installed. Run: pip install reportlab")

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4,
                            leftMargin=2*cm, rightMargin=2*cm,
                            topMargin=2*cm, bottomMargin=2*cm)
    styles = getSampleStyleSheet()

    _register_cn_font(styles)

    content = _content_blocks(params)
    _require_non_empty("content", content)
    elements = []

    for block in content:
        block_type = _block_type(block)
        text = _block_text(block)
        bold = _block_bold(block)
        align = _block_align(block)
        level = _block_level(block)

        align_map = {"left": TA_LEFT, "center": TA_CENTER, "right": TA_RIGHT, "justify": TA_JUSTIFY}
        para_align = align_map.get(align, TA_LEFT)

        if _is_heading(block_type):
            font_size = {1: 22, 2: 18, 3: 16, 4: 14}.get(level, 18)
            style = ParagraphStyle(f"heading{level}", parent=styles["Heading1"],
                                   fontSize=font_size, alignment=para_align,
                                   spaceAfter=12, spaceBefore=18,
                                   fontName="STSong" if bold else "STSong",
                                   textColor=HexColor("#2395bc"))
            elements.append(Paragraph(text, style))
            elements.append(Spacer(1, 6))

        elif _is_paragraph(block_type):
            font_name = "STSong-Bold" if bold else "STSong"
            style = ParagraphStyle("body", parent=styles["Normal"],
                                   fontSize=12, alignment=para_align,
                                   spaceAfter=8, leading=20,
                                   fontName=font_name)
            elements.append(Paragraph(text, style))

        elif _is_table(block_type):
            rows = _normalize_table_data(block)
            if rows:
                col_count = max(len(r) for r in rows) if rows else 1
                table_data = [[_cell_value(c) for c in r] + [""] * (col_count - len(r)) for r in rows]
                table = Table(table_data, colWidths=[4.5*cm] * col_count)
                header_color = HexColor("#2395bc")
                style_cmds = [
                    ("FONTNAME", (0, 0), (-1, 0), "STSong-Bold"),
                    ("TEXTCOLOR", (0, 0), (-1, 0), HexColor("#FFFFFF")),
                    ("BACKGROUND", (0, 0), (-1, 0), header_color),
                    ("FONTNAME", (0, 1), (-1, -1), "STSong"),
                    ("FONTSIZE", (0, 0), (-1, -1), 10),
                    ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                    ("GRID", (0, 0), (-1, -1), 0.5, HexColor("#CCCCCC")),
                    ("TOPPADDING", (0, 0), (-1, -1), 6),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
                ]
                table.setStyle(TableStyle(style_cmds))
                elements.append(table)
                elements.append(Spacer(1, 12))

        elif _is_page_break(block_type):
            elements.append(PageBreak())

    if not elements:
        raise ValueError("content must contain at least one renderable block")

    doc.build(elements)
    buf.seek(0)
    return buf.getvalue()


def _register_cn_font(styles):
    """Register Chinese-supporting fonts for reportlab."""
    try:
        import os

        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont

        font_paths = [
            "/System/Library/Fonts/PingFang.ttc",
            "/System/Library/Fonts/STSong.ttf",
            "/System/Library/Fonts/Supplemental/Songti.ttc",
        ]
        registered = False
        for fp in font_paths:
            if os.path.exists(fp):
                try:
                    pdfmetrics.registerFont(TTFont("STSong", fp))
                    pdfmetrics.registerFont(TTFont("STSong-Bold", fp))
                    registered = True
                    break
                except Exception:
                    continue

        if not registered:
            logger.warning("No Chinese font found for PDF, using Helvetica (Chinese text may be blank)")
    except Exception as exc:
        logger.warning("Font registration failed: %s", exc)
