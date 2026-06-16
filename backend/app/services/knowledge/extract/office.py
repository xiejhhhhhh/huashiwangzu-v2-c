import logging
import os
from app.services.knowledge.extract.types import PageResult

logger = logging.getLogger(__name__)

_DOCLING_AVAILABLE = False
_DOCX_AVAILABLE = False
_PPTX_AVAILABLE = False

try:
    from docling.document_converter import DocumentConverter
    _DOCLING_AVAILABLE = True
except ImportError:
    pass

try:
    import docx
    _DOCX_AVAILABLE = True
except ImportError:
    pass

try:
    import pptx
    _PPTX_AVAILABLE = True
except ImportError:
    pass

_DOCLING_CONVERTER = None


def _get_docling():
    global _DOCLING_CONVERTER
    if _DOCLING_CONVERTER is None:
        _DOCLING_CONVERTER = DocumentConverter()
    return _DOCLING_CONVERTER


def _collect_docx_text(doc) -> tuple[list[str], list[dict]]:
    full_text = []
    layout_blocks = []

    def _add(para):
        text = para.text or ""
        if text.strip():
            full_text.append(text)
            layout_blocks.append({
                "type": "paragraph",
                "style": para.style.name if para.style else "Normal",
                "text": text,
            })

    for para in doc.paragraphs:
        _add(para)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    _add(para)
    for section in doc.sections:
        for h in (section.header, section.first_page_header):
            if h:
                for para in h.paragraphs:
                    _add(para)
        for f in (section.footer, section.first_page_footer):
            if f:
                for para in f.paragraphs:
                    _add(para)

    return full_text, layout_blocks


class OfficeExtractor:

    def extract(self, file_path: str) -> list[PageResult]:
        global _DOCLING_AVAILABLE
        ext = os.path.splitext(file_path)[1].lower()
        try:
            if _DOCLING_AVAILABLE:
                return self._extract_docling(file_path)
        except Exception as e:
            logger.warning("Docling Office extraction failed: %s, falling back", e)
            _DOCLING_AVAILABLE = False

        if ext in (".doc", ".docx"):
            if _DOCX_AVAILABLE:
                return self._extract_docx(file_path)
            raise RuntimeError("python-docx not installed")
        if ext in (".ppt", ".pptx"):
            if _PPTX_AVAILABLE:
                return self._extract_pptx(file_path)
            raise RuntimeError("python-pptx not installed")
        raise ValueError(f"Unsupported office format: {ext}")

    def _extract_docling(self, file_path: str) -> list[PageResult]:
        converter = _get_docling()
        result = converter.convert(file_path)
        pages = []
        for page_num, page in enumerate(result.document.pages.values(), start=1):
            script_text = page.get_text() or ""
            layout_blocks = []
            for item in page.items:
                label = getattr(item, "label", None) or "unknown"
                bbox = getattr(item, "bbox", None)
                text = getattr(item, "text", "") if hasattr(item, "text") else ""
                layout_blocks.append({
                    "type": str(label),
                    "bbox": str(bbox) if bbox else None,
                    "text": text,
                })
            pages.append(PageResult(
                page_num=page_num,
                script_text=script_text,
                layout_data={"blocks": layout_blocks},
            ))
        return pages

    def _extract_docx(self, file_path: str) -> list[PageResult]:
        doc = docx.Document(file_path)
        full_text, layout_blocks = _collect_docx_text(doc)
        body = "\n".join(full_text)
        return [PageResult(page_num=1, script_text=body, layout_data={
            "blocks": layout_blocks,
            "paragraphs": len(layout_blocks),
        })]

    def _extract_pptx(self, file_path: str) -> list[PageResult]:
        prs = pptx.Presentation(file_path)
        pages = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts = []
            layout_blocks = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
                    layout_blocks.append({
                        "type": "text_frame",
                        "shape_name": shape.name,
                        "text": shape.text,
                    })
                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        table_data.append([cell.text for cell in row.cells])
                    layout_blocks.append({
                        "type": "table",
                        "rows": len(shape.table.rows),
                        "cols": len(shape.table.columns),
                        "data": table_data,
                    })
            pages.append(PageResult(
                page_num=slide_num,
                script_text="\n".join(slide_texts),
                layout_data={"blocks": layout_blocks, "shapes": len(layout_blocks)},
            ))
        return pages
