import logging
from pptx import Presentation
from app.core.exceptions import ValidationError

logger = logging.getLogger(__name__)


class PptxService:

    MAX_SLIDES = 200
    MAX_TEXT_LENGTH = 5000

    async def parse(self, file_path: str) -> dict:
        prs = Presentation(file_path)
        if len(prs.slides) > self.MAX_SLIDES:
            raise ValidationError(f"幻灯片数量超过 {self.MAX_SLIDES} 限制")

        content = []
        for idx, slide in enumerate(prs.slides):
            elements = []
            tb_count = 0
            for shape in slide.shapes:
                if shape.has_text_frame:
                    tb_count += 1
                    text = shape.text_frame.text[:self.MAX_TEXT_LENGTH]
                    elements.append({
                        "id": f"s{idx + 1}_tb{tb_count}",
                        "type": "textbox",
                        "content": text,
                        "shape_name": shape.name or "",
                    })
                if shape.shape_type == 13:
                    elements.append({
                        "id": f"s{idx + 1}_img{tb_count}",
                        "type": "image",
                        "resource_id": shape.image.content_type if hasattr(shape, "image") else "",
                    })

            content.append({
                "id": f"s{idx + 1}",
                "index": idx,
                "elements": elements,
            })

        return {
            "manifest": {
                "file_type": "pptx", "version": "1.0.0",
                "slide_count": len(content),
            },
            "content": content,
        }

    async def export(self, file_path: str, json_content: dict) -> None:
        from pptx import Presentation as NewPresentation
        prs = NewPresentation()
        content = json_content.get("content", json_content) if isinstance(json_content, dict) else json_content

        if isinstance(content, list):
            for slide_data in content:
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                elements = slide_data.get("elements", [])
                for elem in elements:
                    if elem.get("type") == "textbox":
                        from pptx.util import Inches, Pt
                        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
                        tf = txBox.text_frame
                        tf.text = elem.get("content", "")[:self.MAX_TEXT_LENGTH]

        prs.save(file_path)

    def preview_patch(self, patch: dict, json_content: dict) -> dict:
        if patch.get("operation_type") not in ("replace_text",):
            raise ValueError("PPTX 补丁仅支持 replace_text 操作类型")
        return {"preview_passed": True, "risk_level": "medium"}
